#!/usr/bin/env python3
"""
Création de la présentation PowerPoint pour la soutenance
Système de Recommandation My Content
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

def create_presentation():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # Couleurs
    RED = RGBColor(192, 0, 0)  # Rouge pour les titres
    BLACK = RGBColor(0, 0, 0)  # Noir pour le texte

    def add_title_slide(title, subtitle):
        """Slide de titre"""
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank

        # Titre
        title_box = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(1))
        title_frame = title_box.text_frame
        title_para = title_frame.paragraphs[0]
        title_para.text = title
        title_para.font.size = Pt(32)
        title_para.font.bold = True
        title_para.font.color.rgb = RED
        title_para.alignment = PP_ALIGN.CENTER

        # Sous-titre
        subtitle_box = slide.shapes.add_textbox(Inches(1), Inches(3.8), Inches(8), Inches(0.8))
        subtitle_frame = subtitle_box.text_frame
        subtitle_para = subtitle_frame.paragraphs[0]
        subtitle_para.text = subtitle
        subtitle_para.font.size = Pt(24)
        subtitle_para.font.color.rgb = BLACK
        subtitle_para.alignment = PP_ALIGN.CENTER

        return slide

    def add_content_slide(title, content_items):
        """Slide avec contenu - Sans puces, avec texte en gras et retraits"""
        slide = prs.slides.add_slide(prs.slide_layouts[6])

        # Titre en rouge
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.8))
        title_frame = title_box.text_frame
        title_para = title_frame.paragraphs[0]
        title_para.text = title
        title_para.font.size = Pt(32)
        title_para.font.bold = True
        title_para.font.color.rgb = RED

        # Contenu en noir - SANS PUCES
        content_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(8.4), Inches(5.5))
        text_frame = content_box.text_frame
        text_frame.word_wrap = True

        for i, item in enumerate(content_items):
            if i > 0:
                text_frame.add_paragraph()
            p = text_frame.paragraphs[i]

            # IMPORTANT: Désactiver explicitement les puces
            p.level = 0

            if isinstance(item, dict):
                # Item avec niveau et style
                p.text = item['text']
                level = item.get('level', 0)

                # Retrait basé sur le niveau (sans puces)
                if level == 0:
                    p.space_before = Pt(8)
                    p.font.size = Pt(item.get('size', 20))
                    p.font.bold = item.get('bold', True)
                elif level == 1:
                    p.space_before = Pt(4)
                    p.left_indent = Inches(0.3)
                    p.font.size = Pt(item.get('size', 17))
                    p.font.bold = item.get('bold', False)
                else:
                    p.left_indent = Inches(0.6)
                    p.font.size = Pt(item.get('size', 16))
                    p.font.bold = False
            else:
                # Item simple (ligne vide ou texte normal)
                p.text = item
                p.font.size = Pt(17)
                p.font.bold = False

            p.font.color.rgb = BLACK

        return slide

    # ==========================================
    # SLIDE 1: Page de titre
    # ==========================================
    add_title_slide(
        "Système de Recommandation Hybride",
        "My Content - Optimisation & Déploiement Cloud"
    )

    # ==========================================
    # SLIDE 2: Contexte et objectifs
    # ==========================================
    add_content_slide(
        "Contexte du Projet",
        [
            {"text": "My Content - Portail d'actualités", "bold": True, "size": 20},
            {"text": "Challenge: Recommander des articles pertinents", "level": 1},
            {"text": "322,897 utilisateurs", "level": 1},
            {"text": "2,872,899 interactions filtrées (règle 30s)", "level": 1},
            {"text": "44,692 articles", "level": 1},
            "",
            {"text": "Objectifs", "bold": True, "size": 20},
            {"text": "Système hybride combinant 3 approches", "level": 1},
            {"text": "Optimisation mémoire (< 30 GB)", "level": 1},
            {"text": "Déploiement cloud production-ready", "level": 1},
            {"text": "API avec latence < 1 seconde", "level": 1},
        ]
    )

    # ==========================================
    # SLIDE 3: Architecture technique
    # ==========================================
    add_content_slide(
        "Architecture Technique",
        [
            {"text": "Pipeline de données", "bold": True, "size": 20},
            {"text": "1. Preprocessing: Matrices sparse (4.4 MB vs 600 GB dense)", "level": 1},
            {"text": "2. Enrichissement: 9 signaux de qualité d'engagement", "level": 1},
            {"text": "3. Matrice pondérée: Weights au lieu de counts", "level": 1},
            {"text": "4. Modèles Lite: 86 MB (réduction 96%)", "level": 1},
            "",
            {"text": "Déploiement", "bold": True, "size": 20},
            {"text": "Azure Functions Consumption Plan", "level": 1},
            {"text": "Runtime: Python 3.11", "level": 1},
            {"text": "Region: France Central", "level": 1},
            {"text": "Endpoint: func-mycontent-reco-1269.azurewebsites.net", "level": 1},
        ]
    )

    # ==========================================
    # SLIDE 4: Système hybride
    # ==========================================
    add_content_slide(
        "Système Hybride - 3 Approches",
        [
            {"text": "1. Content-Based (39%)", "bold": True, "size": 20},
            {"text": "Similarité cosine sur embeddings 250D", "level": 1},
            {"text": "Boost catégories préférées (+10% max)", "level": 1},
            "",
            {"text": "2. Collaborative Filtering (36%)", "bold": True, "size": 20},
            {"text": "Top 50 utilisateurs similaires", "level": 1},
            {"text": "Agrégation pondérée par interaction_weight", "level": 1},
            "",
            {"text": "3. Temporal Scoring (25%)", "bold": True, "size": 20},
            {"text": "Popularité + Temporal decay (half-life 7 jours)", "level": 1},
            {"text": "Crucial pour actualités", "level": 1},
        ]
    )

    # ==========================================
    # SLIDE 5: Optimisation Optuna
    # ==========================================
    add_content_slide(
        "Optimisation Bayésienne des Poids",
        [
            {"text": "Fonction objectif", "bold": True, "size": 20},
            {"text": "Maximiser le temps de lecture (sans temps fantômes < 30s)", "level": 1},
            {"text": "30 trials Optuna TPE avec early stopping", "level": 1},
            "",
            {"text": "Hyperparamètres optimisés: Poids hybrides", "bold": True, "size": 20},
            {"text": "Content-Based: 39% (similarité embeddings)", "level": 1},
            {"text": "Collaborative: 36% (utilisateurs similaires)", "level": 1},
            {"text": "Temporal: 25% (fraîcheur articles)", "level": 1},
            "",
            {"text": "9 signaux de qualité (features fixes)", "bold": True, "size": 20},
            {"text": "Utilisés lors du preprocessing pour pondérer les interactions", "level": 1},
            {"text": "Temps lecture, Clicks, Session, Device, OS, Pays...", "level": 1},
        ]
    )

    # ==========================================
    # SLIDE 6: Optimisation - 9 signaux
    # ==========================================
    add_content_slide(
        "Optimisation: 9 Signaux de Qualité d'Engagement",
        [
            {"text": "Au lieu de simples counts, on calcule un score 0-1:", "size": 18},
            "",
            {"text": "Signaux comportementaux", "bold": True, "size": 20},
            {"text": "Temps de lecture (pondération 40%)", "level": 1},
            {"text": "Nombre de clicks (30%)", "level": 1},
            {"text": "Qualité de session (10%)", "level": 1},
            "",
            {"text": "Signaux contextuels", "bold": True, "size": 20},
            {"text": "Type de device (Desktop > Mobile)", "level": 1},
            {"text": "Système d'exploitation", "level": 1},
            {"text": "Pays, Région, Referrer, Environnement (app/web)", "level": 1},
            "",
            {"text": "Score moyen: 0.353 (range 0.29-0.81)", "bold": True},
        ]
    )

    # ==========================================
    # SLIDE 6: Filtre de qualité 30s
    # ==========================================
    add_content_slide(
        "Filtre de Qualité: Seuil 30 Secondes",
        [
            {"text": "Problématique", "bold": True, "size": 20},
            {"text": "Interactions < 30s = Bruit dans les données", "level": 1},
            {"text": "Clics accidentels, bounces, previews", "level": 1},
            {"text": "Besoin de filtrer pour garder vraies lectures", "level": 1},
            "",
            {"text": "Impact du filtrage", "bold": True, "size": 20},
            {"text": "Avant: 2,872,899 interactions brutes", "level": 1},
            {"text": "Après: 2,420,134 interactions qualité (filtrage 15.7%)", "level": 1},
            {"text": "452,765 interactions < 30s supprimées", "level": 1},
            "",
            {"text": "Résultat: Recommandations basées sur engagement réel", "bold": True},
        ]
    )

    # ==========================================
    # SLIDE 7: Optimisation mémoire
    # ==========================================
    add_content_slide(
        "Défi: Optimisation Mémoire",
        [
            {"text": "Challenge initial", "bold": True, "size": 20},
            {"text": "Versions 1-7: > 40 GB (échec)", "level": 1},
            {"text": "Limite serveur: 30 GB", "level": 1},
            "",
            {"text": "Solution V8 (succès)", "bold": True, "size": 20},
            {"text": "Traitement par batches de 50 fichiers", "level": 1},
            {"text": "Chunks de 5,000 utilisateurs", "level": 1},
            {"text": "Libération mémoire après chaque batch", "level": 1},
            {"text": "Parallélisation contrôlée (12 threads)", "level": 1},
            "",
            {"text": "Résultat final", "bold": True, "size": 20},
            {"text": "4.99 GB / 30 GB (réduction 87.5%)", "level": 1, "bold": True},
            {"text": "Temps: 7 min 48s pour tout le pipeline", "level": 1, "bold": True},
        ]
    )

    # ==========================================
    # SLIDE 8: Pipeline local reproductible
    # ==========================================
    add_content_slide(
        "Pipeline Complet Automatisé",
        [
            {"text": "Exécution en local - 7 étapes", "bold": True, "size": 20},
            {"text": "1. Vérification prérequis (Python, dépendances)", "level": 1},
            {"text": "2. Exploration (364,047 articles, 461 catégories)", "level": 1},
            {"text": "3. Preprocessing (matrices sparse + profils)", "level": 1},
            {"text": "4. Enrichissement (9 signaux, 385 fichiers)", "level": 1},
            {"text": "5. Matrice pondérée (9.0 MB)", "level": 1},
            {"text": "6. Modèles Lite (86 MB)", "level": 1},
            {"text": "7. Validation + Rapport", "level": 1},
            "",
            {"text": "Avantages vs Kaggle", "bold": True, "size": 20},
            {"text": "Reproductible (script bash unique)", "level": 1},
            {"text": "Logs détaillés automatiques", "level": 1},
            {"text": "Pas de limite de temps", "level": 1},
        ]
    )

    # ==========================================
    # SLIDE 9: Résultats techniques
    # ==========================================
    add_content_slide(
        "Résultats Techniques",
        [
            {"text": "Modèles générés", "bold": True, "size": 20},
            {"text": "Matrice: (160,377 × 37,891) = 2,420,134 interactions", "level": 1},
            {"text": "Sparsité: 99.96%", "level": 1},
            {"text": "Profils enrichis: 322,897 utilisateurs", "level": 1},
            "",
            {"text": "Performance API", "bold": True, "size": 20},
            {"text": "Latence moyenne: 650ms (warm)", "level": 1},
            {"text": "Cold start: 715ms", "level": 1},
            {"text": "Objectif à terme: < 200ms", "level": 1},
            "",
            {"text": "Modèles Lite (déploiement)", "bold": True, "size": 20},
            {"text": "86 MB (vs 2.6 GB complets)", "level": 1},
            {"text": "10,000 utilisateurs échantillonnés", "level": 1},
        ]
    )

    # ==========================================
    # SLIDE 10: Impact business
    # ==========================================
    add_content_slide(
        "Impact Business",
        [
            {"text": "Modèle de revenus", "bold": True, "size": 20},
            {"text": "Publicités pop-up: 1 toutes les 3.55 min (médiane)", "level": 1},
            {"text": "CPM: 6€ pour les publicités interstitielles", "level": 1},
            {"text": "322,897 utilisateurs (dataset complet)", "level": 1},
            "",
            {"text": "Résultats (données réelles)", "bold": True, "size": 20},
            {"text": "Sans recommandation: 8,975€", "level": 1},
            {"text": "Avec recommandation: 16,425€", "level": 1},
            {"text": "Gain: +7,450€ (+83% temps passé)", "level": 1, "bold": True},
            "",
            {"text": "ROI Infrastructure", "bold": True, "size": 20},
            {"text": "Coût MVP (Consumption): 122€/an", "level": 1},
            {"text": "ROI net: +6,009% (7,450€ gain - 122€ coût)", "level": 1, "bold": True},
        ]
    )

    # ==========================================
    # SLIDE 11: Démonstration
    # ==========================================
    add_content_slide(
        "Démonstration - Application Streamlit",
        [
            {"text": "Interface interactive développée", "bold": True, "size": 20},
            {"text": "Sélection utilisateur", "level": 1},
            {"text": "4 stratégies prédéfinies (Équilibrée, Personnalisée, Trending, Similaires)", "level": 1},
            {"text": "Mode avancé avec sliders de poids", "level": 1},
            "",
            {"text": "Interprétabilité complète", "bold": True, "size": 20},
            {"text": "Profil utilisateur (articles lus, clicks, temps)", "level": 1},
            {"text": "Catégories préférées vs recommandées", "level": 1},
            {"text": "Noms de catégories (150+ mappés)", "level": 1},
            {"text": "Visualisations Plotly (scores, distribution, temporalité)", "level": 1},
            {"text": "Export CSV/JSON", "level": 1},
        ]
    )

    # ==========================================
    # SLIDE 12: Difficultés surmontées
    # ==========================================
    add_content_slide(
        "Difficultés Techniques Surmontées",
        [
            {"text": "1. Optimisation mémoire", "bold": True, "size": 20},
            {"text": "7 versions échouées avant V8 (4.99 GB)", "level": 1},
            {"text": "Solution: Batching + chunking + libération mémoire", "level": 1},
            "",
            {"text": "2. Latence API", "bold": True, "size": 20},
            {"text": "650ms vs objectif 200ms", "level": 1},
            {"text": "Acceptable pour MVP, optimisations futures identifiées", "level": 1},
            "",
            {"text": "3. Couverture utilisateurs", "bold": True, "size": 20},
            {"text": "Modèles Lite: 10k users seulement", "level": 1},
            {"text": "Solution future: Fallback sur recommandations populaires", "level": 1},
        ]
    )

    # ==========================================
    # SLIDE 13: Livrables
    # ==========================================
    add_content_slide(
        "Livrables du Projet",
        [
            {"text": "Code & Documentation", "bold": True, "size": 20},
            {"text": "Pipeline complet automatisé (run_pipeline_complet.sh)", "level": 1},
            {"text": "API Azure Functions déployée", "level": 1},
            {"text": "Application Streamlit (démonstration)", "level": 1},
            {"text": "Documentation technique exhaustive", "level": 1},
            "",
            {"text": "Modèles", "bold": True, "size": 20},
            {"text": "Modèles complets (2.6 GB)", "level": 1},
            {"text": "Modèles Lite déployés (86 MB)", "level": 1},
            {"text": "Profils enrichis 322k utilisateurs", "level": 1},
            "",
            {"text": "Tests & Validation", "bold": True, "size": 20},
            {"text": "Tests API (7 tests fonctionnels)", "level": 1},
            {"text": "Rapport de tests détaillé", "level": 1},
        ]
    )

    # ==========================================
    # SLIDE 14: Améliorations futures
    # ==========================================
    add_content_slide(
        "Perspectives d'Amélioration",
        [
            {"text": "Court terme", "bold": True, "size": 20},
            {"text": "Optimisation latence (profiling + cache)", "level": 1},
            {"text": "Azure Premium Plan (< 200ms)", "level": 1},
            {"text": "Fallback recommandations populaires", "level": 1},
            "",
            {"text": "Moyen terme", "bold": True, "size": 20},
            {"text": "Features supplémentaires (geo, session complète)", "level": 1},
            {"text": "A/B testing en production", "level": 1},
            {"text": "Feedback utilisateurs (clicks réels)", "level": 1},
            "",
            {"text": "Long terme", "bold": True, "size": 20},
            {"text": "Learning-to-Rank (LightGBM)", "level": 1},
            {"text": "Deep Learning (BERT4Rec, Transformers)", "level": 1},
            {"text": "Personnalisation temps réel", "level": 1},
        ]
    )

    # ==========================================
    # SLIDE 15: Conclusion
    # ==========================================
    add_content_slide(
        "Conclusion",
        [
            {"text": "Réalisations clés", "bold": True, "size": 20},
            {"text": "Système hybride optimisé (39/36/25)", "level": 1},
            {"text": "Optimisation mémoire réussie (87.5% réduction)", "level": 1},
            {"text": "Déploiement Azure fonctionnel", "level": 1},
            {"text": "ROI exceptionnel: +6,009%", "level": 1},
            {"text": "Pipeline reproductible et automatisé", "level": 1},
            "",
            {"text": "Points forts", "bold": True, "size": 20},
            {"text": "Optimisation Optuna: poids hybrides (39/36/25)", "level": 1},
            {"text": "Target: temps de lecture (sans temps fantômes)", "level": 1},
            {"text": "Interprétabilité complète", "level": 1},
            {"text": "Documentation exhaustive + Application démo", "level": 1},
            "",
            {"text": "Prêt pour production My Content", "bold": True, "size": 22},
        ]
    )

    # ==========================================
    # SLIDE 16: Questions
    # ==========================================
    add_title_slide(
        "Questions ?",
        "Merci de votre attention"
    )

    # Sauvegarder
    output_file = "/home/ser/Bureau/P10_reco_new/PRESENTATION_SOUTENANCE.pptx"
    prs.save(output_file)
    print(f"✅ Présentation créée: {output_file}")
    print(f"   Nombre de slides: {len(prs.slides)}")
    return output_file

if __name__ == "__main__":
    create_presentation()
