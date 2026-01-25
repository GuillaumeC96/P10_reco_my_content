#!/usr/bin/env python3
"""
Script de génération automatique de la présentation PowerPoint
Projet P10 - My Content - Système de Recommandation d'Articles

Usage:
    python3 generer_presentation.py

Output:
    LIVRABLES_PROJET10_Cassez_Guillaume_012026/3_Cassez_Guillaume_3_presentation_122024/Cassez_Guillaume_3_presentation_122024.pptx
"""

import re
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# Couleurs du thème (bleu professionnel)
COLOR_TITLE = RGBColor(31, 78, 120)      # Bleu foncé
COLOR_SUBTITLE = RGBColor(68, 114, 196)  # Bleu moyen
COLOR_TEXT = RGBColor(64, 64, 64)        # Gris foncé
COLOR_ACCENT = RGBColor(0, 176, 80)      # Vert pour les ✅
COLOR_WARNING = RGBColor(255, 153, 0)    # Orange pour les ⚠️
COLOR_ERROR = RGBColor(192, 0, 0)        # Rouge pour les ❌

def parse_markdown_content(content):
    """Parse le contenu Markdown et extrait les slides."""
    slides = []
    current_slide = None

    # Split par les slides
    slide_sections = content.split('\n## SLIDE')

    for section in slide_sections:
        if not section.strip():
            continue

        lines = section.strip().split('\n')

        # Première ligne contient le numéro et titre
        if lines:
            first_line = lines[0]
            # Extraire le titre (enlever le numéro)
            title_match = re.match(r'(\d+)\s*[-–]\s*(.+)', first_line)
            if title_match:
                slide_num = int(title_match.group(1))
                title = title_match.group(2).strip()

                # Le reste du contenu
                content_lines = lines[1:]
                content = '\n'.join(content_lines).strip()

                slides.append({
                    'number': slide_num,
                    'title': title,
                    'content': content
                })

    return slides

def clean_text(text):
    """Nettoie le texte des marqueurs Markdown."""
    # Enlever les ** pour le gras
    text = re.sub(r'\*\*(.+?)\*\*', r'\1', text)
    # Enlever les ` pour le code
    text = re.sub(r'`(.+?)`', r'\1', text)
    # Enlever les liens []()
    text = re.sub(r'\[(.+?)\]\(.+?\)', r'\1', text)
    return text

def add_title_slide(prs, slide_info):
    """Ajoute la slide de titre."""
    slide = prs.slides.add_slide(prs.slide_layouts[0])  # Layout titre

    title = slide.shapes.title
    subtitle = slide.placeholders[1]

    title.text = "My Content"
    title.text_frame.paragraphs[0].font.size = Pt(54)
    title.text_frame.paragraphs[0].font.bold = True
    title.text_frame.paragraphs[0].font.color.rgb = COLOR_TITLE

    subtitle.text = "Système de Recommandation d'Articles\n\nEncourager la lecture par des recommandations pertinentes\n\nGuillaume Cassez - CTO & Co-fondateur\nDécembre 2024"
    subtitle.text_frame.paragraphs[0].font.size = Pt(20)
    subtitle.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

def add_content_slide(prs, slide_info):
    """Ajoute une slide de contenu."""
    slide = prs.slides.add_slide(prs.slide_layouts[1])  # Layout titre + contenu

    title = slide.shapes.title
    title.text = slide_info['title']
    title.text_frame.paragraphs[0].font.size = Pt(32)
    title.text_frame.paragraphs[0].font.bold = True
    title.text_frame.paragraphs[0].font.color.rgb = COLOR_TITLE

    # Zone de contenu
    content_box = slide.placeholders[1]
    tf = content_box.text_frame
    tf.clear()

    # Parser le contenu
    content = slide_info['content']
    lines = content.split('\n')

    for line in lines:
        line = line.strip()
        if not line or line.startswith('---') or line.startswith('**'):
            continue

        # Créer un paragraphe
        p = tf.add_paragraph()

        # Détecter le niveau (nombre de tirets ou espaces)
        indent_level = 0
        clean_line = line

        if line.startswith('- '):
            indent_level = 0
            clean_line = line[2:]
        elif line.startswith('  - '):
            indent_level = 1
            clean_line = line[4:]
        elif line.startswith('    '):
            indent_level = 1
            clean_line = line[4:]
        elif line.startswith('○ '):
            indent_level = 1
            clean_line = line[2:]

        # Nettoyer le texte
        clean_line = clean_text(clean_line)

        p.text = clean_line
        p.level = indent_level
        p.font.size = Pt(18 - indent_level * 2)
        p.font.color.rgb = COLOR_TEXT

        # Colorer les symboles spéciaux
        if '✅' in clean_line or '✓' in clean_line:
            p.font.color.rgb = COLOR_ACCENT
        elif '❌' in clean_line or '✗' in clean_line:
            p.font.color.rgb = COLOR_ERROR
        elif '⚠️' in clean_line:
            p.font.color.rgb = COLOR_WARNING

def add_section_slide(prs, title):
    """Ajoute une slide de séparation de section."""
    slide = prs.slides.add_slide(prs.slide_layouts[5])  # Layout vide

    # Ajouter un titre centré
    left = Inches(1)
    top = Inches(3)
    width = Inches(8)
    height = Inches(1.5)

    textbox = slide.shapes.add_textbox(left, top, width, height)
    tf = textbox.text_frame
    tf.text = title

    p = tf.paragraphs[0]
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = COLOR_TITLE
    p.alignment = PP_ALIGN.CENTER

def add_architecture_slide(prs, slide_info):
    """Ajoute une slide avec un schéma d'architecture."""
    slide = prs.slides.add_slide(prs.slide_layouts[1])

    title = slide.shapes.title
    title.text = slide_info['title']
    title.text_frame.paragraphs[0].font.size = Pt(32)
    title.text_frame.paragraphs[0].font.bold = True
    title.text_frame.paragraphs[0].font.color.rgb = COLOR_TITLE

    # Zone de contenu avec le schéma ASCII
    content_box = slide.placeholders[1]
    tf = content_box.text_frame
    tf.clear()

    # Extraire le schéma ASCII du contenu
    content = slide_info['content']

    # Si contient un bloc de code avec le schéma
    if '```' in content:
        schema_match = re.search(r'```(.*?)```', content, re.DOTALL)
        if schema_match:
            schema = schema_match.group(1).strip()
            p = tf.add_paragraph()
            p.text = schema
            p.font.name = 'Courier New'
            p.font.size = Pt(14)
            p.font.color.rgb = COLOR_TEXT
    else:
        # Sinon ajouter le contenu normal
        add_content_slide(prs, slide_info)

def add_table_slide(prs, slide_info):
    """Ajoute une slide avec un tableau comparatif."""
    slide = prs.slides.add_slide(prs.slide_layouts[5])  # Layout vide

    title = slide.shapes.title
    title.text = slide_info['title']
    title.text_frame.paragraphs[0].font.size = Pt(32)
    title.text_frame.paragraphs[0].font.bold = True
    title.text_frame.paragraphs[0].font.color.rgb = COLOR_TITLE

    # Extraire les lignes du tableau du contenu
    content = slide_info['content']
    if '|' in content:
        # Parser le tableau Markdown
        table_lines = [line.strip() for line in content.split('\n') if '|' in line]
        if len(table_lines) >= 2:
            # Créer un tableau
            rows = len(table_lines) - 1  # -1 pour la ligne de séparation
            cols = len(table_lines[0].split('|')) - 2  # -2 pour les | aux extrémités

            left = Inches(0.5)
            top = Inches(2)
            width = Inches(9)
            height = Inches(4)

            table = slide.shapes.add_table(rows, cols, left, top, width, height).table

            # Remplir les cellules
            row_idx = 0
            for line in table_lines:
                if '---' in line:  # Skip separator line
                    continue
                cells = [cell.strip() for cell in line.split('|')[1:-1]]
                for col_idx, cell_text in enumerate(cells):
                    if row_idx < rows and col_idx < cols:
                        cell = table.rows[row_idx].cells[col_idx]
                        cell.text = clean_text(cell_text)
                        cell.text_frame.paragraphs[0].font.size = Pt(14)
                        if row_idx == 0:  # Header
                            cell.text_frame.paragraphs[0].font.bold = True
                            cell.fill.solid()
                            cell.fill.fore_color.rgb = COLOR_SUBTITLE
                row_idx += 1
    else:
        # Pas de tableau, ajouter comme slide normale
        add_content_slide(prs, slide_info)

def generate_presentation():
    """Génère la présentation PowerPoint complète."""
    print("🎨 Génération de la présentation PowerPoint...")

    # Chemins
    base_dir = Path(__file__).parent
    content_file = base_dir / "CONTENU_PRESENTATION.md"
    output_dir = base_dir / "LIVRABLES_PROJET10_Cassez_Guillaume_012026" / "3_Cassez_Guillaume_3_presentation_122024"
    output_file = output_dir / "Cassez_Guillaume_3_presentation_122024.pptx"

    # Vérifier que le fichier de contenu existe
    if not content_file.exists():
        print(f"❌ Erreur: Fichier {content_file} introuvable")
        return False

    # Créer le dossier de sortie si nécessaire
    output_dir.mkdir(parents=True, exist_ok=True)

    # Lire le contenu
    print(f"📖 Lecture du contenu depuis {content_file}...")
    with open(content_file, 'r', encoding='utf-8') as f:
        content = f.read()

    # Parser les slides
    print("🔍 Parsing du contenu...")
    slides = parse_markdown_content(content)
    print(f"   Trouvé {len(slides)} slides")

    # Créer la présentation
    print("📊 Création de la présentation...")
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # Sections pour organiser
    sections = {
        1: "Introduction",
        4: "Dataset & Application",
        7: "Analyse des Modèles",
        13: "Architecture & Système",
        19: "Architecture Cible",
        25: "Conclusion"
    }

    # Ajouter les slides
    for i, slide_info in enumerate(slides, 1):
        slide_num = slide_info['number']

        print(f"   Slide {slide_num}: {slide_info['title']}")

        # Ajouter slide de section si nécessaire
        if slide_num in sections:
            section_title = sections[slide_num]
            if slide_num > 1:  # Pas de section avant la première slide
                print(f"   → Section: {section_title}")
                add_section_slide(prs, section_title)

        # Ajouter la slide selon son type
        if slide_num == 1:
            add_title_slide(prs, slide_info)
        elif slide_num in [12, 20]:  # Slides avec schémas d'architecture
            add_architecture_slide(prs, slide_info)
        elif slide_num == 11:  # Slide avec tableau comparatif
            add_table_slide(prs, slide_info)
        else:
            add_content_slide(prs, slide_info)

    # Sauvegarder
    print(f"💾 Sauvegarde de la présentation...")
    prs.save(str(output_file))
    print(f"✅ Présentation créée avec succès!")
    print(f"   Fichier: {output_file}")
    print(f"   Nombre de slides: {len(prs.slides)}")

    return True

def main():
    """Point d'entrée principal."""
    print("=" * 70)
    print("GÉNÉRATION PRÉSENTATION POWERPOINT - PROJET P10")
    print("My Content - Système de Recommandation d'Articles")
    print("=" * 70)
    print()

    try:
        success = generate_presentation()

        if success:
            print()
            print("=" * 70)
            print("✅ GÉNÉRATION TERMINÉE AVEC SUCCÈS!")
            print("=" * 70)
            print()
            print("📌 PROCHAINES ÉTAPES:")
            print("   1. Ouvrir le fichier .pptx")
            print("   2. Vérifier le contenu et la mise en forme")
            print("   3. Ajuster les schémas si nécessaire")
            print("   4. Ajouter des images/icônes si désiré")
            print("   5. Exporter en PDF pour soumission (optionnel)")
            print()
            print("💡 Le fichier PowerPoint est prêt pour la soutenance!")
            print()
        else:
            print()
            print("❌ ÉCHEC DE LA GÉNÉRATION")
            print("   Vérifiez les messages d'erreur ci-dessus")
            print()
            return 1

    except Exception as e:
        print()
        print(f"❌ ERREUR: {e}")
        import traceback
        traceback.print_exc()
        return 1

    return 0

if __name__ == "__main__":
    exit(main())
